from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    # Use a blank slide layout (usually index 6)
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)

    # 1. Set Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 230, 220)  # Light beige

    # Dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    width_in = slide_width.inches
    height_in = slide_height.inches

    # Colors
    col_dark_green = RGBColor(28, 78, 50)
    col_med_green = RGBColor(60, 145, 90)
    col_light_green_arrow = RGBColor(160, 200, 170)
    col_text_q = RGBColor(28, 78, 50)
    col_white = RGBColor(255, 255, 255)

    # Data Structure
    # (Question, [Right Side Blocks Data])
    # Block Data: (Text, Type, Count) 
    # Type: 'roof', 'rect', 'split', '5-split', '4-split'
    
    rows = [
        ("你是谁？", {"text": ["品牌定位"], "type": "roof"}),
        ("你主张什么？", {"text": ["品牌主张"], "type": "rect"}),
        ("你的风格是什么？", {"text": ["品牌调性"], "type": "rect"}),
        ("你的口号是什么？", {"text": ["品牌口号"], "type": "rect"}),
        ("你的目标人群是谁？", {"text": ["核心目标人群"], "type": "rect"}),
        ("你能给TA带来什么利益？", {"text": ["功能性利益点", "情感性利益点"], "type": "split_2"}),
        ("你凭什么被TA相信？", {"text": ["RTB1", "RTB2", "RTB3", "RTB4", "RTB5"], "type": "split_5"}),
        ("你都有什么产品？", {"text": ["产品系列1", "产品系列2", "产品系列3", "产品系列4"], "type": "split_4"})
    ]

    # Layout Config
    margin_top = Inches(0.5)
    margin_bottom = Inches(0.5)
    
    content_height = height_in - 1.0 # 0.5 margin top/bottom
    row_height = content_height / len(rows)
    
    # Left Column (Questions)
    left_col_x = Inches(0.5)
    left_col_w = Inches(3.0)
    
    # Arrow Column
    arrow_x = left_col_x + left_col_w
    arrow_w = Inches(1.0)
    
    # Right Column (Pyramid)
    right_col_x = arrow_x + arrow_w
    right_col_w = width_in - right_col_x - Inches(0.5)

    gap = Inches(0.05) # Gap between blocks

    for i, (question, block_data) in enumerate(rows):
        y_pos = margin_top + (i * row_height)
        
        # 1. Draw Question
        txBox = slide.shapes.add_textbox(left_col_x, y_pos + (row_height/4), left_col_w, row_height/2)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = question
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = col_text_q
        p.alignment = PP_ALIGN.LEFT

        # 2. Draw Arrow
        # Using a chevron or similar, or just a custom shape. 
        # Let's use a simple arrow shape pointing left.
        arrow_shape = slide.shapes.add_shape(
            MSO_SHAPE.LEFT_ARROW, 
            arrow_x + Inches(0.1), 
            y_pos + (row_height * 0.3), 
            arrow_w - Inches(0.2), 
            row_height * 0.4
        )
        arrow_shape.fill.solid()
        arrow_shape.fill.fore_color.rgb = col_light_green_arrow
        arrow_shape.line.fill.background() # No line

        # 3. Draw Right Blocks
        b_type = block_data["type"]
        texts = block_data["text"]
        
        # Determine block width based on count
        if b_type == "roof":
            # Triangle top
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE, # Or CHEVRON, but Triangle is better for "roof"
                right_col_x, 
                y_pos, 
                right_col_w, 
                row_height - gap
            )
            # Flatten the triangle a bit by making it wide? 
            # Actually ISOSCELES_TRIANGLE points up by default.
            
            shape.fill.solid()
            shape.fill.fore_color.rgb = col_dark_green # Top one is darker
            shape.line.fill.background()
            
            shape.text_frame.text = texts[0]
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            shape.text_frame.paragraphs[0].font.color.rgb = col_white
            shape.text_frame.paragraphs[0].font.bold = True
            
        elif b_type == "rect":
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                right_col_x, 
                y_pos, 
                right_col_w, 
                row_height - gap
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = col_med_green
            shape.line.fill.background()
            
            shape.text_frame.text = texts[0]
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            shape.text_frame.paragraphs[0].font.color.rgb = col_white
            shape.text_frame.paragraphs[0].font.bold = True

        elif "split" in b_type:
            count = len(texts)
            block_w = (right_col_w - (gap * (count - 1))) / count
            
            for j, text in enumerate(texts):
                bx = right_col_x + (j * (block_w + gap))
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 
                    bx, 
                    y_pos, 
                    block_w, 
                    row_height - gap
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = col_med_green
                shape.line.fill.background()
                
                shape.text_frame.text = text
                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                shape.text_frame.paragraphs[0].font.color.rgb = col_white
                shape.text_frame.paragraphs[0].font.bold = True
                shape.text_frame.paragraphs[0].font.size = Pt(12 if count > 2 else 14)

    prs.save('brand_pyramid.pptx')
    print("Presentation saved as brand_pyramid.pptx")

if __name__ == "__main__":
    create_presentation()
