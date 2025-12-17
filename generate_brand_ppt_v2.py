from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def create_presentation_v2():
    prs = Presentation()
    
    # 1. Force 16:9 Aspect Ratio (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)

    # 2. Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 230, 220)  # Light beige

    # 3. Configuration
    rows = [
        ("你是谁？", {"text": ["品牌定位"], "type": "roof"}),
        ("你主张什么？", {"text": ["品牌主张"], "type": "rect"}),
        ("你的风格是什么？", {"text": ["品牌调性"], "type": "rect"}),
        ("你的口号是什么？", {"text": ["品牌口号"], "type": "rect"}),
        ("你的目标人群是谁？", {"text": ["核心目标人群"], "type": "rect"}),
        ("你能给TA带来什么利益？", {"text": ["功能性利益点", "情感性利益点"], "type": "split_2"}),
        ("你凭什么被TA相信？", {"text": ["RTB1", "RTB2", "RTB3", "RTB4", "RTB5"], "type": "split_5"}),
        ("你都有什么产品？", {"text": ["产品系列1", "产品系列2", "产品系列3", "产品系列4", "产品系列5"], "type": "split_5_products"}) 
        # Note: User image had 4 product series in first request, but 5 in second OCR. 
        # The second OCR says "Product Series 1...5". I will use 5.
    ]

    # Colors
    col_dark_green = RGBColor(28, 78, 50)
    col_med_green = RGBColor(60, 145, 90)
    col_light_green_arrow = RGBColor(160, 200, 170)
    col_text_q = RGBColor(28, 78, 50)
    col_white = RGBColor(255, 255, 255)

    # Layout Calculations
    margin_y = Inches(0.5)
    margin_x = Inches(0.5)
    
    total_height = prs.slide_height.inches - (2 * margin_y.inches)
    row_h_in = total_height / len(rows)
    row_height = Inches(row_h_in)
    
    # Column Widths
    # Left (Questions): 2.5 inches
    # Arrow: 0.8 inches
    # Right (Pyramid): Rest
    
    col_q_x = margin_x
    col_q_w = Inches(2.5)
    
    col_arrow_x = col_q_x + col_q_w
    col_arrow_w = Inches(0.8)
    
    col_pyramid_x = col_arrow_x + col_arrow_w
    col_pyramid_w = prs.slide_width - col_pyramid_x - margin_x # Remaining width

    gap = Inches(0.05)

    for i, (question, block_data) in enumerate(rows):
        # Calculate Y position for this row
        y = margin_y + (i * row_height)
        
        # --- 1. Question Text ---
        # Vertically center the text box in the row
        # Text box height slightly smaller than row height to avoid touch
        tb_h = row_height * 0.8
        tb_y = y + (row_height - tb_h) / 2
        
        txBox = slide.shapes.add_textbox(col_q_x, tb_y, col_q_w, tb_h)
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE # Critical for vertical centering
        p = tf.paragraphs[0]
        p.text = question
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = col_text_q
        p.alignment = PP_ALIGN.LEFT

        # --- 2. Arrow ---
        arrow_h = row_height * 0.4
        arrow_y = y + (row_height - arrow_h) / 2
        # Arrow points left. 
        # We can use LEFT_ARROW shape.
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.LEFT_ARROW,
            col_arrow_x, arrow_y, col_arrow_w - Inches(0.1), arrow_h
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = col_light_green_arrow
        arrow.line.fill.background()

        # --- 3. Pyramid Blocks ---
        b_type = block_data["type"]
        texts = block_data["text"]
        
        block_h = row_height - gap
        block_y = y # Align to top of row (with gap handled by height reduction usually, or just tight packing)
        # Actually better to center vertically if there's space, but for pyramid usually they stack tightly.
        # Let's use tight stacking.
        
        # Determine shape type and widths
        if b_type == "roof":
            # Single block, Triangle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                col_pyramid_x, block_y, col_pyramid_w, block_h
            )
            color = col_dark_green
            
        elif b_type == "rect":
            # Single block, Rect
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                col_pyramid_x, block_y, col_pyramid_w, block_h
            )
            color = col_med_green
            
        elif "split" in b_type:
            # Multiple blocks
            count = len(texts)
            # Calculate width per block
            # Total width = (count * w) + ((count-1) * gap)
            # w = (Total - (count-1)*gap) / count
            
            total_gap_w = gap * (count - 1)
            single_w = (col_pyramid_w - total_gap_w) / count
            
            for j, txt in enumerate(texts):
                bx = col_pyramid_x + (j * (single_w + gap))
                
                rect = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    bx, block_y, single_w, block_h
                )
                rect.fill.solid()
                rect.fill.fore_color.rgb = col_med_green
                rect.line.fill.background()
                
                # Text Frame
                tf = rect.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.text = txt
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = col_white
                p.font.bold = True
                
                # Dynamic font size
                if len(txt) > 6 or count > 4:
                    p.font.size = Pt(10)
                elif count > 2:
                    p.font.size = Pt(12)
                else:
                    p.font.size = Pt(14)
            
            continue # Skip the single-shape logic below

        # Common logic for single shapes (roof/rect)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE # Center vertically
        # Note: Triangle vertical anchor might be tricky, usually BOTTOM or MIDDLE works.
        # For Isosceles Triangle, the text box is the bounding box. 
        # Text might overlap the top peak if font is big. 
        # We can add a margin top to text frame or just keep font reasonable.
        
        if b_type == "roof":
            tf.margin_top = Inches(0.1) # Push down slightly from the sharp point
            
        p = tf.paragraphs[0]
        p.text = texts[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = col_white
        p.font.bold = True
        p.font.size = Pt(16)

    output_path = 'brand_pyramid_v2.pptx'
    prs.save(output_path)
    print(f"Presentation saved as {output_path}")

if __name__ == "__main__":
    create_presentation_v2()
